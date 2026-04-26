#!/usr/bin/env python3
"""Build R3 PPT from R2 template with images, proper formatting."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

BASE = "/Users/prashantkumar/Desktop/Luna_Aura"
TMPL = os.path.join(BASE, "archive/legacy_2026_cleanup/Documents /PPT R2_B904.pptx")
OUT = os.path.join(BASE, "LunaAura_Review3_Final.pptx")
IMG = os.path.join(BASE, "ppt_images")
FIGS = os.path.join(BASE, "paper_assets")

prs = Presentation(TMPL)

# Remove all existing slides
while len(prs.slides._sldIdLst) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get(qn('r:id'))
    if rId:
        try: prs.part.drop_rel(rId)
        except: pass
    prs.slides._sldIdLst.remove(sldId)

TB = prs.slide_layouts[3]   # Title & Bullets
TI = prs.slide_layouts[0]   # Title
BL = prs.slide_layouts[16]  # Blank
PH = prs.slide_layouts[15]  # Photo

def set_tr(run, size=20):
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)

def tb_slide(title, bullets, bullet_size=18):
    """Title & Bullets slide with proper text sizing."""
    s = prs.slides.add_slide(TB)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0 or ph.top < Emu(2500000):
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs: set_tr(r, 36)
        else:
            tf = ph.text_frame; tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                p.space_after = Pt(4)
                p.space_before = Pt(2)
                for r in p.runs: set_tr(r, bullet_size)
    return s

def img_slide(title, img_path, caption=""):
    """Full-image slide with title."""
    s = prs.slides.add_slide(BL)
    # Add title textbox at top
    tx = s.shapes.add_textbox(Emu(1200000), Emu(300000), Emu(22000000), Emu(1200000))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    for r in p.runs: set_tr(r, 36)
    # Add image centered
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Emu(3500000), Emu(1800000), Emu(17000000))
    if caption:
        cx = s.shapes.add_textbox(Emu(1200000), Emu(12500000), Emu(22000000), Emu(800000))
        cp = cx.text_frame.paragraphs[0]
        cp.text = caption
        for r in cp.runs: set_tr(r, 14)
    return s

def dual_img_slide(title, img1, img2, cap1="", cap2=""):
    """Slide with two images side by side."""
    s = prs.slides.add_slide(BL)
    tx = s.shapes.add_textbox(Emu(1200000), Emu(300000), Emu(22000000), Emu(1200000))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    for r in p.runs: set_tr(r, 36)
    if os.path.exists(img1):
        s.shapes.add_picture(img1, Emu(500000), Emu(2000000), Emu(11500000))
    if os.path.exists(img2):
        s.shapes.add_picture(img2, Emu(12500000), Emu(2000000), Emu(11500000))
    return s

def title_slide(title, lines):
    s = prs.slides.add_slide(TI)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs: set_tr(r, 40)
        else:
            tf = ph.text_frame; tf.clear()
            for i, l in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = l
                for r in p.runs: set_tr(r, 20)
    return s

# ===================== SLIDES =====================

# 1 TITLE
title_slide(
    "LunaAura: An Explainable Machine Learning Framework\nfor Hormonal and Emotional Health Analysis",
    ["Eshita Rai  |  er8137@srmist.edu.in",
     "Guide: Mrs. P. Sivakamasundari",
     "Department of Computing Technologies  |  SRM IST, Chennai",
     "Review 3  |  Batch B904"]
)

# 2 ABSTRACT
tb_slide("Abstract", [
    "Emotional wellbeing depends on sleep, stress, hormonal cycles, and daily habits",
    "— but most wellness apps treat these in isolation with no depth.",
    "",
    "LunaAura combines a transparent risk formula with calibrated ML,",
    "quantile regression, and SHAP explainability in one platform.",
    "",
    "It tracks six daily inputs, computes personalized risk and wellness scores,",
    "and delivers actionable insights through a longitudinal dashboard.",
    "",
    "Tested on 105 synthetic users (470 daily records) + 40-day case study.",
    "Positioned as a research prototype — no clinical diagnostic claim."
])

# 3 MOTIVATION
tb_slide("Motivation", [
    "970 million people worldwide live with a mental health condition (WHO, 2022)",
    "",
    "University students face compounding pressures — irregular sleep,",
    "academic stress, financial strain — often unnoticed until crisis.",
    "",
    "Current wellness apps fall short:",
    "  • Assess each factor in isolation (sleep alone, stress alone)",
    "  • Ignore the menstrual cycle despite decades of clinical evidence",
    "  • Never explain why a score was generated",
    "",
    "→ Need for a smarter, multi-dimensional, explainable, cycle-aware system"
])

# 4 PROBLEM STATEMENT
tb_slide("Problem Statement", [
    "Emotional wellbeing fluctuates with lifestyle and hormonal context.",
    "",
    "Existing platforms have critical gaps:",
    "  • No integrated behavioral + hormonal analysis",
    "  • Limited predictive capability — mostly descriptive tools",
    "  • No explainable ML insights — users get numbers without reasons",
    "  • No gender-adaptive logic for equitable analytics",
    "",
    "→ A personalized, transparent, cycle-aware framework is needed",
    "   that provides actionable wellness insights with explainability."
])

# 5 OBJECTIVES
tb_slide("Research Objectives", [
    "1. Design a transparent 6-factor composite risk scoring formula",
    "2. Incorporate menstrual cycle phase modifiers from chronobiology",
    "3. Build gender-adaptive weight normalization for all users",
    "4. Train calibrated ML models on PHQ-9 data",
    "5. Implement SHAP-based per-instance explainability",
    "6. Develop a full-stack interactive dashboard with history tracking",
    "7. Validate on a reproducible 105-user synthetic cohort"
])

# 6 LITERATURE SURVEY
tb_slide("Literature Survey", [
    "Mental Health Screening:",
    "  • PHQ-9 (Kroenke, 2001) — gold standard depression measure",
    "  • Arefin & de Roux (2020) — ensemble classifiers for depressive text",
    "",
    "Digital Phenotyping:",
    "  • Saeb et al. (2015) — phone sensors predict PHQ-9 variation",
    "  • Torous et al. (2019, 2021) — data quality and interpretability gaps",
    "",
    "Explainable AI:",
    "  • SHAP (Lundberg & Lee, 2017) — game-theoretic feature attribution",
    "  • EmotionSense (Rachuri, 2010) — contextual feedback boosts engagement",
    "",
    "Cycle-Aware Health:",
    "  • Baker & Driver (2007) — sleep disruption in luteal phase",
    "  • Kiesner (2012) — individual variation in cycle-mood effects"
])

# 7 RESEARCH GAP
tb_slide("Research Gap", [
    "No existing platform combines all of these:",
    "",
    "  ✗ Multi-factor behavioral risk scoring",
    "  ✗ Menstrual cycle phase awareness in scoring models",
    "  ✗ Gender-adaptive weight normalization",
    "  ✗ Calibrated probabilistic ML with uncertainty estimation",
    "  ✗ Per-instance SHAP explainability",
    "  ✗ Longitudinal tracking with rolling temporal features",
    "",
    "LunaAura integrates all six in a single framework.",
    "The novelty is in the integration, not a new algorithm."
])

# 8 SDG
tb_slide("SDG Alignment", [
    "SDG 3 — Good Health & Well-Being",
    "  • Mental health awareness through accessible, personalized analytics",
    "",
    "SDG 5 — Gender Equality",
    "  • Cycle-phase-aware scoring designed specifically for women's health",
    "  • Gender-adaptive design ensures equal depth for all users",
    "",
    "SDG 9 — Industry, Innovation & Infrastructure",
    "  • Novel hybrid deterministic + ML + SHAP architecture",
    "",
    "SDG 10 — Reduced Inequalities",
    "  • Free, browser-based — no expensive hardware needed"
])

# 9 PHASE 1 — ORIGINAL VERSION
tb_slide("Phase 1 — What LunaAura Started As", [
    "The initial prototype had:",
    "",
    "  • Basic HTML input forms for wellness data",
    "  • Simple heuristic scoring — no ML pipeline",
    "  • Static Chart.js visualizations with hardcoded ranges",
    "  • Login attempt without persistent user profiles",
    "  • Single-page dashboard, limited interactivity",
    "  • No longitudinal tracking — each session was independent",
    "  • No cycle-aware scoring or gender logic",
    "  • Generic, copy-paste recommendations for every user"
])

# 10 PHASE 2 — PROBLEMS
tb_slide("Phase 2 — Problems We Identified", [
    "Honest engineering reflection:",
    "",
    "  • Static charts did not update with real user data",
    "  • No personalization — same output for every user",
    "  • No backend persistence — data lost on page refresh",
    "  • No historical analytics — no trend tracking over time",
    "  • Generic outputs — recommendations felt shallow",
    "  • No explainability — scores with no context",
    "  • Gender logic broken — Male selection removed all inputs",
    "  • No research or analytics module for population insights"
])

# 11 PHASE 3 — IMPROVEMENTS
tb_slide("Phase 3 — Improvements Implemented", [
    "Complete system overhaul:",
    "",
    "  • SQLite backend with persistent user profiles",
    "  • 105 pseudo-users across three risk strata",
    "  • 470 daily log entries with realistic distributions",
    "  • Dynamic charts that update in real-time",
    "  • 30-40 day longitudinal history tracking",
    "  • Gender-adaptive dashboard (hides only cycle elements for males)",
    "  • Analytics endpoint for population summaries",
    "  • Model Insights page with feature importance",
    "  • Research Mode for deeper exploration",
    "  • Professional UI with responsive layout"
])

# 12 ARCHITECTURE DIAGRAM
img_slide("System Architecture",
    os.path.join(IMG, "architecture.png"),
    "Three-tier: Frontend (HTML/JS/Chart.js) → Flask REST API (7 endpoints) → SQLite Database")

# 13 ML PIPELINE
img_slide("Machine Learning Pipeline",
    os.path.join(IMG, "ml_pipeline.png"),
    "Data Collection → KNN Linking → Feature Engineering → 3 Model Branches → Prediction API")

# 14 INPUTS
tb_slide("Behavioral Inputs Collected Daily", [
    "Six daily behavioral signals:",
    "",
    "  • Sleep Duration — hours slept (8h optimal)",
    "  • Perceived Stress — ordinal 1-10",
    "  • Physical Activity — minutes/day (60 min target)",
    "  • Anxiety Level — ordinal 1-10",
    "  • Water Intake — litres/day (2.5L target)",
    "  • Menstrual Cycle Day — day 1-28 (female only)",
    "",
    "Plus demographics: Age, Gender, Height, Weight, Cycle Length"
])

# 15 RISK FORMULA
tb_slide("Composite Risk Scoring Formula", [
    "R = Stress(35%) + SleepDeficit(25%) + Anxiety(20%)",
    "    + ActivityDeficit(10%) + Hydration(5%) + Baseline(5) + Phase(±5-8)",
    "",
    "Cycle Phase Modifiers (female only):",
    "  Menstrual (Day 1-5):   +5     Follicular (Day 6-13):  -5",
    "  Ovulatory (Day 14-16): -8     Luteal (Day 17-28):     +8",
    "",
    "Male/Other: phase = 0, sleep weight rises 25% → 30%",
    "",
    "Severity Tiers:",
    "  Low (R < 35)  |  Moderate (35 ≤ R < 65)  |  High (R ≥ 65)",
    "",
    "Design: no single factor alone triggers High severity"
])

# 16 WELLNESS
tb_slide("Wellness Score Estimation", [
    "W = Σ αⱼ × x̂ⱼ   (each x̂ⱼ normalized to [0,100])",
    "",
    "Female (8 factors):",
    "  Sleep 25% | Stress 20% | Activity 15% | Anxiety 15%",
    "  Cycle 10% | Water 5% | Age 5% | Stability 5%",
    "",
    "Male/Other (6 factors — cycle removed, weights redistributed):",
    "  Sleep 30% | Stress 25% | Activity 20%",
    "  Anxiety 15% | Water 5% | Age 5%",
    "",
    "Normalization: Sleep penalizes deviation from 8h optimum",
    "Stress inversely mapped; Activity scaled to 60-min target"
])

# 17 ML MODELS
tb_slide("Machine Learning Models", [
    "1. Calibrated Classifier",
    "   HistGradientBoosting (100 rounds, 31 leaves, η=0.05)",
    "   Platt scaling via 5-fold stratified CV",
    "   Target: PHQ-9 referral flag (n₀=5,054 / n₁=4,494)",
    "",
    "2. Quantile Regressors (×3)",
    "   HistGradientBoosting (150 rounds, pinball loss)",
    "   q = 0.10, 0.50, 0.90 → PHQ-9 confidence intervals",
    "",
    "3. SHAP Explanation Proxy",
    "   RandomForest (100 trees, depth 5)",
    "   TreeExplainer → exact Shapley values per feature",
    "   Top 3 attributions shown as plain-language insights",
    "",
    "Training corpus: 9,548 records × 28 columns"
])

# 18 NOVELTY
tb_slide("Technical Novelty", [
    "What sets LunaAura apart (honestly stated):",
    "",
    "✦ Hybrid Architecture — deterministic (auditable) + ML (probabilistic)",
    "",
    "✦ Cycle-Aware Scoring — first prototype with hormonal phase modifiers",
    "",
    "✦ Gender-Adaptive Normalization — re-normalized weights, not removed features",
    "",
    "✦ Longitudinal Features — 3-day rolling means capture behavioral momentum",
    "",
    "✦ Instance-Level Explainability — SHAP insight cards on dashboard",
    "",
    "We do not claim a new ML algorithm.",
    "Our novelty is the integration framework."
])

# 19 DATABASE
tb_slide("Database Design — SQLite", [
    "Database: lunaaura.db",
    "",
    "Table: users (12 fields)",
    "  id, username, password_hash, age, gender, height, weight,",
    "  cycle_length, sleep_target, cohort_group, source, timestamps",
    "",
    "Table: user_history (12 fields per entry)",
    "  id, user_id (FK), date, sleep_duration, stress_level,",
    "  mood_score, cycle_day, phase, activity, wellness_score,",
    "  predicted_risk, anxiety_level, water_liters",
    "",
    "Cohort: 105 users (86F / 14M / 5O) | 470 daily logs",
    "Ages 18-55 (mean 36.5) | Eshita: 85 entries, 40+ days"
])

# 20 LOGIN SCREENSHOT
img_slide("User Interface — Login & Authentication",
    os.path.join(IMG, "login_modal.png"),
    "Secure login with demo user profiles for testing")

# 21 DAILY LOG SCREENSHOT
img_slide("User Interface — Daily Behavioral Logging",
    os.path.join(IMG, "daily_log.png"),
    "Slider-based input for sleep, stress, activity, anxiety, water, cycle day")

# 22 DASHBOARD SCREENSHOT
img_slide("User Interface — Dashboard Results",
    os.path.join(IMG, "dashboard_results.png"),
    "Real-time chart rendering after data submission")

# 23 ANALYTICS SCREENSHOT
img_slide("Analytics & Model Insights",
    os.path.join(IMG, "analytics_dashboard.png"),
    "Population-level summaries and model metadata")

# 24 RESULTS — ROC + FEATURE IMPORTANCE
dual_img_slide("Results — Model Performance",
    os.path.join(FIGS, "15_roc_curve.png"),
    os.path.join(FIGS, "18_feature_importance.png"))

# 25 RESULTS — CORRELATION + SENSITIVITY
dual_img_slide("Results — Validation Charts",
    os.path.join(FIGS, "14_correlation_heatmap.png"),
    os.path.join(FIGS, "21_sensitivity_risk_vs_stress.png"))

# 26 RESULTS INTERPRETATION
tb_slide("Results Interpretation", [
    "Classifier: AUROC ≈ 0.567 | Brier Score = 0.245",
    "  Modest but honestly reported — reflects cross-dataset noise",
    "",
    "Feature Importance: Age > Stress > Cycle Day > Sleep",
    "  Rolling temporal features contribute meaningfully",
    "",
    "Sensitivity Analysis:",
    "  Risk vs Stress is monotonic and clinically plausible",
    "  Moderate threshold crossed at stress ≈ 5.5",
    "  High tier requires multi-factor adversity",
    "",
    "Cycle Phase Impact: 16-point spread (Ovulatory -8 to Luteal +8)",
    "  Intentionally modest — modulates, does not dominate"
])

# 27 LIMITATIONS & VALIDATION
tb_slide("Limitations & Validation", [
    "Limitations:",
    "  • Synthetic cohort (seed=42) — not clinically collected",
    "  • Formula weights are expert heuristics, not data-optimized",
    "  • ML AUROC reflects heterogeneous data bridging",
    "  • Cycle modifiers are population averages",
    "  • System produces Wellness Estimates — NOT clinical diagnoses",
    "",
    "Validation Performed:",
    "  ✓ Sensitivity analysis — monotonic, directionally correct",
    "  ✓ Correlation heatmap validates formula weighting",
    "  ✓ All 7 API endpoints tested",
    "  ✓ Frontend works for Female, Male, Other profiles",
    "  ✓ Database integrity: 105 users, 470 logs, FK constraints",
    "  ✓ Charts validated against formula outputs"
])

# 28 COST & FEASIBILITY
tb_slide("Economic & Financial Feasibility", [
    "Development Cost: ₹0 — student-built, open-source tools",
    "",
    "Technology Stack (all free):",
    "  • Python + Flask — backend",
    "  • SQLite — zero-config database",
    "  • scikit-learn — ML models",
    "  • SHAP — explainability",
    "  • Chart.js — visualization",
    "  • HTML/CSS/JS — frontend",
    "",
    "Deployment: runs on any laptop | scalable to cloud",
    "SQLite → PostgreSQL migration for production",
    "",
    "Affordable, reproducible, student-feasible solution"
])

# 29 ETHICS
tb_slide("Ethics & Privacy", [
    "1. No Clinical Claims",
    "   Output labeled as Behavioral Wellness Estimate",
    "",
    "2. Privacy-First Design",
    "   All data stored locally — no external transmission",
    "   User consent required for real-world deployment",
    "",
    "3. Explainable AI",
    "   SHAP provides transparent per-instance attribution",
    "   Users understand why they received their score",
    "",
    "4. Bias Awareness",
    "   Synthetic data disclosed; population-level modifiers acknowledged",
    "   Gender-adaptive design prevents analytical disparity",
    "",
    "5. SDG 3 Alignment — Good Health and Well-Being"
])

# 30 METHODOLOGY FLOW
tb_slide("Methodology — End-to-End Pipeline", [
    "User Input (6 signals + demographics)",
    "    ↓",
    "Flask API /predict receives JSON payload",
    "    ↓",
    "Deterministic Risk Formula → R ∈ [0,100]",
    "    ↓",
    "ML Inference:",
    "  • Calibrated Classifier → probability",
    "  • Quantile Regressors → PHQ-9 confidence interval",
    "  • SHAP Proxy → top-3 feature attributions",
    "    ↓",
    "Recommendation Engine → priority-tagged suggestions",
    "    ↓",
    "SQLite Logging → entry persisted",
    "    ↓",
    "Dashboard Rendering → all modules update (< 2 seconds)"
])

# 31 PAPER
tb_slide("Research Paper Publication", [
    "Title: LunaAura: An Explainable ML Framework for",
    "       Hormonal and Emotional Health Analysis",
    "",
    "Authors: Eshita Rai, Mrs. P. Sivakamasundari",
    "Format: IEEE conference paper (6 pages, two-column)",
    "",
    "Contents:",
    "  • 11 formal equations",
    "  • 5 publication-ready figures",
    "  • 2 data tables",
    "  • 18 scholarly references",
    "",
    "Status: Complete, formatted, submission-ready"
])

# 32 CONCLUSION
tb_slide("Conclusion", [
    "Problem: Wellness tools lack multi-factor scoring,",
    "         cycle awareness, and explainability",
    "",
    "Solution: LunaAura — hybrid deterministic + ML framework",
    "",
    "Objectives Achieved:",
    "  ✓ Transparent 6-factor risk formula with auditable weights",
    "  ✓ Cycle-phase modifiers from published chronobiology",
    "  ✓ Gender-adaptive weight normalization",
    "  ✓ Calibrated classifier + quantile regression + SHAP",
    "  ✓ Full-stack dashboard with longitudinal tracking",
    "  ✓ Validated on 105-user cohort (470 logs, 40-day study)",
    "",
    "Outputs: AUROC 0.567 | Brier 0.245 | 9,548 training records"
])

# 33 FUTURE SCOPE
tb_slide("Future Scope", [
    "Planned extensions:",
    "",
    "  • Clinical validation with real datasets under IRB approval",
    "  • Wearable integration — HRV, actigraphy, sleep staging",
    "  • Federated learning for privacy-preserving training",
    "  • Bayesian weight optimization on clinical outcomes",
    "  • Native mobile app for continuous monitoring",
    "  • Individualized cycle-length modeling (replace 28-day fixed)",
    "",
    "Source code and model artifacts available on GitHub",
    "for independent replication and extension."
])

# 34 THANK YOU
title_slide("Thank You", [
    "Questions Welcome",
    "",
    "Eshita Rai  |  er8137@srmist.edu.in",
    "Guide: Mrs. P. Sivakamasundari",
    "SRM Institute of Science and Technology, Chennai",
    "Batch B904"])

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
