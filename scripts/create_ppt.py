#!/usr/bin/env python3
"""Generate LunaAura Review 3 PPT from the R2 template."""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(BASE, "archive/legacy_2026_cleanup/Documents /PPT R2_B904.pptx")
OUTPUT = os.path.join(BASE, "LunaAura_Review3_Final.pptx")
ASSETS = os.path.join(BASE, "paper_assets")

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x00, 0xB8, 0x94)
ORANGE = RGBColor(0xFD, 0x79, 0x72)

def load_template():
    prs = Presentation(TEMPLATE)
    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId_attr = list(prs.slides._sldIdLst[0].attrib.values())
            # Just delete the XML element
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
            continue
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return prs

def new_prs():
    """Create fresh presentation with template dimensions."""
    prs = Presentation()
    prs.slide_width = Emu(24384000)
    prs.slide_height = Emu(13716000)
    return prs

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def add_title_shape(slide, text, top=Inches(0.8), left=Inches(1.5), width=Inches(23), size=44, color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_body(slide, bullets, top=Inches(2.8), left=Inches(1.5), width=Inches(21), size=24, color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(9))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(14)
        p.level = 0
    return txBox

def add_accent_bar(slide, top=Inches(2.3)):
    shape = slide.shapes.add_shape(1, Inches(1.5), top, Inches(4), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_image_safe(slide, path, left, top, width):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width)

def make_slide(prs, title, bullets, img_path=None):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    add_title_shape(slide, title)
    add_accent_bar(slide)
    if img_path and os.path.exists(img_path):
        add_body(slide, bullets, width=Inches(12))
        add_image_safe(slide, img_path, Inches(14.5), Inches(3), Inches(10))
    else:
        add_body(slide, bullets)
    return slide

# ====================== BUILD SLIDES ======================
prs = new_prs()

# --- Slide 1: Title ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_shape(s, "LunaAura", top=Inches(3), left=Inches(1.5), size=72, color=ACCENT)
add_title_shape(s, "An Explainable Machine Learning Framework for\nHormonal and Emotional Health Analysis", top=Inches(5), left=Inches(1.5), size=32, color=WHITE, bold=False)
add_body(s, [
    "Eshita Rai  |  er8137@srmist.edu.in",
    "Guide: Mrs. P. Sivakamasundari",
    "Department of Computing Technologies  |  SRM Institute of Science and Technology, Chennai",
    "",
    "Review 3 Presentation  |  Batch B904"
], top=Inches(8), size=22, color=LIGHT)

# --- Slide 2: Why This Project? ---
make_slide(prs, "Why This Project?", [
    "• 970 million people worldwide live with a mental health condition (WHO, 2022)",
    "• University students face compounding stress, sleep disruption, and social isolation",
    "• Existing wellness apps are shallow: step counters, mood emojis, generic tips",
    "• No commercial tool combines multi-factor scoring + hormonal awareness + explainability",
    "• The menstrual cycle reshapes sleep, stress reactivity, and mood — yet trackers ignore it",
    "",
    "→ LunaAura fills this gap with an explainable, cycle-aware, personalized wellness framework"
])

# --- Slide 3: Objectives ---
make_slide(prs, "Research Objectives", [
    "1. Design a transparent composite risk scoring formula integrating six behavioral factors",
    "2. Incorporate menstrual cycle phase modifiers grounded in chronobiological evidence",
    "3. Build a gender-adaptive weight normalization scheme for equitable analytics",
    "4. Train calibrated ML models (HistGradientBoosting + Quantile Regression) on PHQ-9 data",
    "5. Implement SHAP-based per-instance explainability for user-facing insights",
    "6. Develop a full-stack interactive dashboard with longitudinal tracking capability",
    "7. Validate the system on a reproducible 105-user synthetic cohort (470 daily records)"
])

# --- Slide 4: Original Version ---
make_slide(prs, "Phase 1 — Original LunaAura Version", [
    "What the project started with:",
    "",
    "• Basic wellness concept with standard user input forms",
    "• Simple heuristic-based scoring (no ML pipeline)",
    "• Static Chart.js visualizations with hardcoded data ranges",
    "• Initial login/signup attempt without persistent user profiles",
    "• Single-page dashboard with limited analytics depth",
    "• No longitudinal tracking — each session was independent",
    "• No cycle-aware scoring or gender-adaptive logic",
    "• Generic recommendations identical for all users"
])

# --- Slide 5: Problems Found ---
make_slide(prs, "Phase 2 — Problems Identified in Original Version", [
    "Engineering Reflection — what needed fixing:",
    "",
    "• Static charts: visualizations did not update with real user data",
    "• Weak personalization: same output regardless of user profile",
    "• No backend persistence: data lost on page refresh",
    "• No historical analytics: impossible to track trends over days/weeks",
    "• Generic outputs: recommendations were copy-paste templates",
    "• No explainability: users received scores without understanding why",
    "• Gender logic broken: selecting Male removed all inputs, not just cycle fields",
    "• No research or analytics module for population-level insights"
])

# --- Slide 6: Final Architecture ---
make_slide(prs, "Phase 3 — Final System Architecture", [
    "Three-tier modular design:",
    "",
    "Frontend (HTML / JavaScript / Chart.js)",
    "    ↓  REST API calls  ↓",
    "Flask Backend (port 5001) — 7 endpoints",
    "    /health  /signup  /login  /profile  /predict  /analytics  /insights",
    "    ↓",
    "Four Subsystems:",
    "    • Hybrid Scoring Engine (Deterministic + ML)",
    "    • SHAP TreeExplainer Module",
    "    • Rules-Based Recommendation Generator",
    "    • Chart Simulation Engine",
    "    ↓",
    "SQLite Database (lunaaura.db) — users + user_history tables"
])

# --- Slide 7: Inputs Collected ---
make_slide(prs, "Behavioral Inputs Collected Daily", [
    "Six behavioral dimensions tracked per session:",
    "",
    "  Signal                  Scale           Baseline",
    "  ─────────────────────────────────────────────────",
    "  Sleep Duration          0–12 hours      8h optimal",
    "  Perceived Stress        1–10 ordinal    Mid-range",
    "  Physical Activity       0–120 minutes   60 min target",
    "  Anxiety Level           1–10 ordinal    Mid-range",
    "  Water Intake            0–5 litres      2.5L target",
    "  Menstrual Cycle Day     1–28            Phase-dependent",
    "",
    "+ Demographics: Age, Gender, Height, Weight, Cycle Length, Sleep Target"
])

# --- Slide 8: Risk & Wellness Formulas ---
make_slide(prs, "Wellness & Risk Scoring Formulas", [
    "Composite Risk Score R (0–100):",
    "  R = Stress(35%) + SleepDeficit(25%) + Anxiety(20%) + ActivityDeficit(10%)",
    "      + HydrationDeficit(5%) + TrendBaseline(5) + PhaseModifier(±5–8)",
    "",
    "Cycle Phase Modifiers (Female only):",
    "  Menstrual (Day 1–5):  +5    Follicular (Day 6–13): −5",
    "  Ovulatory (Day 14–16): −8   Luteal (Day 17–28):    +8",
    "",
    "Severity: Low (<35)  |  Moderate (35–64)  |  High (≥65)",
    "",
    "Wellness Score W (0–100):",
    "  Female: Sleep 25% + Stress 20% + Activity 15% + Anxiety 15%",
    "          + Cycle 10% + Water 5% + Age 5% + Stability 5%",
    "  Male:   Sleep 30% + Stress 25% + Activity 20% + Anxiety 15%",
    "          + Water 5% + Age 5%  (re-normalized, no cycle)"
])

# --- Slide 9: ML Models ---
make_slide(prs, "Machine Learning Models Used", [
    "1. Calibrated Classifier",
    "   • HistGradientBoostingClassifier (100 rounds, 31 leaves, η=0.05)",
    "   • Sigmoid calibration via 5-fold stratified CV (Platt scaling)",
    "   • Target: PHQ-9-derived referral flag (n₀=5,054 / n₁=4,494)",
    "",
    "2. Quantile Regressors (×3)",
    "   • HistGradientBoostingRegressor (150 rounds, pinball loss)",
    "   • Quantiles: q = 0.10, 0.50, 0.90 → confidence intervals on PHQ-9",
    "",
    "3. SHAP Explanation Proxy",
    "   • RandomForestClassifier (100 trees, depth 5)",
    "   • TreeExplainer computes exact Shapley values per instance",
    "   • Top 3 features surfaced as natural-language insights",
    "",
    "Training data: 9,548 records × 28 columns (master_dataset.csv)"
])

# --- Slide 10: Technical Novelty ---
make_slide(prs, "Technical Novelty — What Makes LunaAura Different", [
    "✦ Hybrid Architecture: deterministic formula (transparent) + ML (probabilistic)",
    "   → Users get auditable scores AND uncertainty-aware estimates",
    "",
    "✦ Cycle-Aware Contextual Intelligence",
    "   → First wellness prototype to integrate hormonal phase modifiers into scoring",
    "",
    "✦ Gender-Adaptive Weight Normalization",
    "   → Male/non-binary users get re-normalized weights, not reduced analytics",
    "",
    "✦ Longitudinal User-Specific Monitoring",
    "   → 30–40 day history tracking with rolling feature engineering",
    "",
    "✦ Instance-Level Explainability",
    "   → SHAP-powered plain-language insight cards on the dashboard",
    "",
    "⚠ We do NOT claim invention of a new ML algorithm",
    "  Our novelty is in the integration and application framework"
])

# --- Slide 11: Database Design ---
make_slide(prs, "Database Design — SQLite Schema", [
    "Table: users (12 fields)",
    "  id | username | password_hash | age | gender | height_cm | weight_kg",
    "  cycle_length | sleep_target | created_at | updated_at | cohort_group | source",
    "",
    "Table: user_history (12 fields)",
    "  id | user_id (FK) | date | sleep_duration | stress_level | mood_score",
    "  cycle_day | phase | activity | wellness_score | predicted_risk",
    "  anxiety_level | water_liters",
    "",
    "Cohort Statistics:",
    "  • 105 registered users (86 Female, 14 Male, 5 Other)",
    "  • 470 daily log entries",
    "  • Ages 18–55 (mean 36.5)",
    "  • 1 longitudinal super-user: Eshita (85 entries, 40+ days)"
])

# --- Slide 12: Dashboard Screenshots ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_shape(s, "Dashboard — Live System Screenshots")
add_accent_bar(s)
add_body(s, [
    "Interactive modules rendered in real-time after signal injection:",
    "",
    "• Wellness Score Ring with animated progress indicator",
    "• 7-Day Risk Momentum Heatmap (color-coded severity)",
    "• 30-Day Trend Projections: wellness, stress, sleep trajectories",
    "• Risk Distribution Doughnut Chart (Low / Moderate / High)",
    "• Factor Contribution Breakdown with signed impact values",
    "• Menstrual Phase Influence Panel (female users only — hidden for male)",
    "• Personalized Insight Cards powered by SHAP attributions",
    "• Recommendation Engine: priority-tagged behavioral suggestions",
    "",
    "All modules update within 2 seconds of data submission"
], width=Inches(21))

# --- Slide 13: Analytics Module ---
make_slide(prs, "Population Analytics Module", [
    "The /analytics endpoint provides cohort-level summaries:",
    "",
    "• Average wellness score across all users",
    "• Average stress level distribution",
    "• Total user count and log entry count",
    "• Risk category breakdown (Low / Moderate / High proportions)",
    "",
    "The /insights endpoint exposes model metadata:",
    "• Model type and training parameters",
    "• Feature list and importance rankings",
    "• Brier calibration score (0.245)",
    "• AUROC metric (~0.567)",
    "",
    "Purpose: enables researcher-facing exploration beyond individual dashboards"
])

# --- Slide 14: Results / Charts ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_shape(s, "Key Results — Model Evaluation & Validation")
add_accent_bar(s)
add_body(s, [
    "Classifier: AUROC ≈ 0.567 | Brier Score = 0.245",
    "  → Modest but honestly reported; reflects cross-dataset noise",
    "",
    "Feature Importance: Age > Stress > Cycle Day > Sleep > Activity",
    "  → Rolling temporal features contribute non-trivially",
    "",
    "Sensitivity Analysis: Risk vs Stress is monotonic and clinically plausible",
    "  → Moderate threshold crossed at stress ≈ 5.5",
    "",
    "Cycle Phase Impact: 16-point spread (Ovulatory −8 to Luteal +8)",
    "  → Intentionally modest — modulates, does not dominate"
], width=Inches(12))
roc = os.path.join(ASSETS, "15_roc_curve.png")
feat = os.path.join(ASSETS, "18_feature_importance.png")
add_image_safe(s, roc, Inches(14), Inches(3), Inches(5))
add_image_safe(s, feat, Inches(19.5), Inches(3), Inches(5))

# --- Slide 15: Limitations + Validation ---
make_slide(prs, "Limitations & Validation  [Rubric: 5.6.1 — 3 marks]", [
    "Limitations (honestly disclosed):",
    "  • Evaluation cohort is synthetic (seed=42), not clinically collected",
    "  • Composite formula uses expert-derived weights, not data-optimized",
    "  • ML classifier's modest AUROC reflects heterogeneous data bridging",
    "  • Cycle-phase modifiers are population averages, not individualized",
    "  • System produces Behavioral Wellness Estimates — NOT clinical diagnoses",
    "",
    "Validation Performed:",
    "  ✓ Sensitivity analysis confirms monotonic, directionally correct formula responses",
    "  ✓ Correlation heatmap validates formula weighting scheme (non-circular)",
    "  ✓ All 7 API endpoints tested with valid and edge-case inputs",
    "  ✓ Frontend modules render correctly for Female, Male, and Other profiles",
    "  ✓ Database integrity verified: 105 users, 470 logs, FK constraints intact",
    "  ✓ Charts validated against deterministic formula outputs"
])

# --- Slide 16: Cost & Feasibility ---
make_slide(prs, "Economic & Financial Feasibility  [Rubric: 11.5.1 — 3 marks]", [
    "Development Cost: ₹0 (open-source stack, student-built)",
    "",
    "Technology Stack — All Free/Open-Source:",
    "  • Python 3 + Flask (backend)            — Free",
    "  • SQLite (database)                      — Free, zero-config",
    "  • scikit-learn (ML models)               — Free",
    "  • SHAP (explainability)                  — Free",
    "  • Chart.js (visualization)               — Free",
    "  • HTML/CSS/JS (frontend)                 — Free",
    "",
    "Deployment Options:",
    "  • Local: runs on any laptop (Flask + SQLite, no server needed)",
    "  • Cloud: easily deployable to Heroku / AWS / GCP free tiers",
    "  • Scaling: SQLite → PostgreSQL migration for production use",
    "",
    "Conclusion: Affordable, reproducible, student-feasible solution"
])

# --- Slide 17: Ethics & Privacy ---
make_slide(prs, "Ethics & Privacy  [Rubric: 8.4.2 — 4 marks]", [
    "Ethical Principles Applied:",
    "",
    "1. No Clinical Claims",
    "   → System explicitly labeled as Behavioral Wellness Estimate, never diagnosis",
    "",
    "2. Privacy-First Design",
    "   → All data stored locally in SQLite; no external transmission",
    "   → User consent required for any future real-world deployment",
    "",
    "3. Explainable AI",
    "   → SHAP provides transparent, per-instance feature attribution",
    "   → Users understand WHY they received a particular score",
    "",
    "4. Bias Awareness",
    "   → Synthetic data disclosed; population-level modifiers acknowledged",
    "   → Gender-adaptive design prevents analytical disparity",
    "",
    "5. SDG 3 Alignment — Good Health and Well-Being"
])

# --- Slide 18: Methodology Summary ---
make_slide(prs, "Methodology Flow  [Rubric: 4.5.1 — 4 marks]", [
    "End-to-End Pipeline:",
    "",
    "  User Input (6 behavioral signals + demographics)",
    "      ↓",
    "  Flask API /predict endpoint receives JSON payload",
    "      ↓",
    "  Deterministic Risk Formula (Eq. 1–7) → Risk Score R ∈ [0,100]",
    "      ↓",
    "  ML Inference: Calibrated Classifier → probability",
    "                Quantile Regressors → PHQ-9 confidence interval",
    "                SHAP Proxy → top-3 feature attributions",
    "      ↓",
    "  Recommendation Engine: threshold-based behavioral suggestions",
    "      ↓",
    "  SQLite Logging: entry persisted to user_history table",
    "      ↓",
    "  Dashboard Rendering: Chart.js updates all modules in <2 seconds"
])

# --- Slide 19: Conclusion ---
make_slide(prs, "Conclusions  [Rubric: 2.8.4 — 3 marks]", [
    "Problem Identified:",
    "  → Existing wellness tools lack multi-factor scoring, cycle awareness, explainability",
    "",
    "Solution Delivered:",
    "  → LunaAura: hybrid deterministic + ML framework with SHAP explainability",
    "",
    "Objectives Achieved:",
    "  ✓ Transparent 6-factor risk formula with auditable weights",
    "  ✓ Cycle-phase modifiers grounded in published chronobiology",
    "  ✓ Gender-adaptive weight normalization (not feature removal)",
    "  ✓ Calibrated classifier + quantile regression + SHAP attribution",
    "  ✓ Full-stack interactive dashboard with longitudinal tracking",
    "  ✓ Validated on 105-user cohort (470 logs, 40-day case study)",
    "",
    "Measurable Outputs: AUROC 0.567 | Brier 0.245 | 9,548 training records"
])

# --- Slide 20: Future Scope + Thank You ---
make_slide(prs, "Future Scope & Thank You", [
    "Planned Extensions:",
    "  • Clinical validation with real longitudinal datasets under IRB approval",
    "  • Wearable integration: heart rate variability, actigraphy, sleep staging",
    "  • Federated learning for privacy-preserving distributed model training",
    "  • Bayesian weight optimization conditioned on clinical outcomes",
    "  • Native mobile app for continuous daily monitoring",
    "  • Individualized cycle-length modeling (replace fixed 28-day assumption)",
    "",
    "Paper Status:",
    "  • IEEE-formatted research paper completed and ready for submission",
    "  • Full codebase available on GitHub for independent replication",
    "",
    "",
    "                    Thank You — Questions Welcome",
    "",
    "         Eshita Rai  |  er8137@srmist.edu.in  |  Batch B904"
])

# Save
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
